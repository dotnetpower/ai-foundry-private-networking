#!/usr/bin/env python3
"""
테스트 문서 생성 및 AI Search 인덱싱 스크립트
- DOCX, PPTX 테스트 문서 생성
- Azure Blob Storage 업로드
- AI Search 인덱스 구성
"""

import os
import json
import sys
from pathlib import Path
from datetime import datetime

# 필요한 패키지 확인 및 설치 안내
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    print("python-docx 패키지가 필요합니다: pip install python-docx")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
except ImportError:
    print("python-pptx 패키지가 필요합니다: pip install python-pptx")
    sys.exit(1)

try:
    from azure.storage.blob import BlobServiceClient, BlobClient
    from azure.identity import DefaultAzureCredential
except ImportError:
    print("azure-storage-blob, azure-identity 패키지가 필요합니다")
    print("pip install azure-storage-blob azure-identity")
    sys.exit(1)

try:
    from azure.search.documents import SearchClient
    from azure.search.documents.indexes import SearchIndexClient
    from azure.search.documents.indexes.models import (
        SearchIndex,
        SearchField,
        SearchFieldDataType,
        SimpleField,
        SearchableField,
        VectorSearch,
        HnswAlgorithmConfiguration,
        VectorSearchProfile,
        SemanticConfiguration,
        SemanticField,
        SemanticPrioritizedFields,
        SemanticSearch,
    )
except ImportError:
    print("azure-search-documents 패키지가 필요합니다")
    print("pip install azure-search-documents")
    sys.exit(1)


# =============================================================================
# 설정
# =============================================================================
STORAGE_ACCOUNT = "staifoundry20260128"
CONTAINER_NAME = "documents"
SEARCH_SERVICE = "srch-aifoundry-7kkykgt6"
INDEX_NAME = "aifoundry-docs-index"
OUTPUT_DIR = Path("./test_documents")


# =============================================================================
# 테스트 문서 콘텐츠
# =============================================================================
DOCUMENTS = {
    "AI_Foundry_소개.docx": {
        "title": "Azure AI Foundry 플랫폼 소개",
        "category": "플랫폼",
        "sections": [
            {
                "heading": "Azure AI Foundry란?",
                "content": """Azure AI Foundry는 Microsoft Azure에서 제공하는 통합 AI 개발 플랫폼입니다. 
이 플랫폼은 기업이 AI 솔루션을 빠르고 안전하게 구축, 배포, 관리할 수 있도록 설계되었습니다.

주요 특징:
- 통합 AI 개발 환경: 모델 개발부터 배포까지 단일 플랫폼에서 관리
- 엔터프라이즈 보안: Private Endpoint, VNet 통합으로 데이터 보호
- 다양한 AI 모델 지원: GPT-4o, DALL-E, Whisper 등 최신 모델 제공
- RAG (Retrieval-Augmented Generation) 지원: AI Search와 연동한 지식 기반 AI"""
            },
            {
                "heading": "AI Hub와 Project 구조",
                "content": """Azure AI Foundry는 Hub와 Project의 계층 구조로 구성됩니다.

AI Hub (허브):
- 중앙 관리 허브로 여러 프로젝트를 관리
- 공유 리소스 (Storage, Key Vault, ACR) 관리
- 네트워크 설정 및 보안 정책 적용

AI Project (프로젝트):
- 개별 AI 솔루션 개발 단위
- Hub에서 리소스를 상속받아 사용
- 팀별 또는 애플리케이션별로 분리"""
            },
            {
                "heading": "프라이빗 네트워킹",
                "content": """엔터프라이즈 환경에서는 프라이빗 네트워킹이 필수입니다.

Private Endpoint 구성:
- AI Hub Private Endpoint: AI Foundry 서비스 접근
- Storage Private Endpoint: Blob, File 스토리지 접근
- Key Vault Private Endpoint: 비밀 키 관리
- OpenAI Private Endpoint: GPT 모델 API 접근
- AI Search Private Endpoint: 검색 서비스 접근

Private DNS Zone:
- privatelink.api.azureml.ms
- privatelink.blob.core.windows.net
- privatelink.vaultcore.azure.net
- privatelink.openai.azure.com
- privatelink.search.windows.net"""
            }
        ]
    },
    "RAG_패턴_가이드.docx": {
        "title": "RAG (Retrieval-Augmented Generation) 패턴 구현 가이드",
        "category": "개발 가이드",
        "sections": [
            {
                "heading": "RAG 패턴이란?",
                "content": """RAG(Retrieval-Augmented Generation)는 대규모 언어 모델(LLM)의 응답을 개선하기 위해 
외부 지식 소스를 활용하는 AI 패턴입니다.

작동 원리:
1. 사용자 질문 수신
2. 벡터 검색으로 관련 문서 검색
3. 검색된 문서를 LLM 프롬프트에 포함
4. 컨텍스트 기반 응답 생성

장점:
- 최신 정보 반영 가능
- 환각(Hallucination) 감소
- 도메인 특화 지식 활용
- 출처 명시 가능"""
            },
            {
                "heading": "Azure AI Search 연동",
                "content": """Azure AI Search는 RAG 패턴의 핵심 구성 요소입니다.

인덱스 구성:
- 문서 콘텐츠 필드 (Searchable)
- 메타데이터 필드 (Filterable)
- 벡터 필드 (1536 차원, Ada-002 Embedding)

Semantic Search 설정:
- prioritizedFields로 중요 필드 지정
- Semantic ranker로 결과 재순위화

하이브리드 검색:
- 키워드 검색 + 벡터 검색 결합
- 정확도와 관련성 모두 향상"""
            },
            {
                "heading": "AI Foundry Playground에서 RAG 사용",
                "content": """AI Foundry Playground에서 RAG를 사용하는 방법:

1. Playground 접속
   - AI Foundry Portal > Project > Playground

2. Add your data 설정
   - 'Add your data' 버튼 클릭
   - Azure AI Search 선택
   - 인덱스 연결 (aifoundry-docs-index)

3. 검색 설정
   - Search type: Hybrid (권장)
   - Semantic search: 활성화
   - Top-k: 5 (검색 결과 개수)

4. 테스트
   - 문서 관련 질문 입력
   - 출처 포함 응답 확인"""
            }
        ]
    },
    "보안_가이드.docx": {
        "title": "Azure AI Foundry 보안 모범 사례",
        "category": "보안",
        "sections": [
            {
                "heading": "제로 트러스트 원칙",
                "content": """Azure AI Foundry는 제로 트러스트 보안 원칙을 따릅니다.

핵심 원칙:
- Never Trust, Always Verify: 모든 접근 요청 검증
- Assume Breach: 침해 가정 하에 설계
- Least Privilege: 최소 권한 원칙

구현 방법:
- Managed Identity: 서비스 간 인증
- RBAC: 역할 기반 접근 제어
- Private Endpoint: 네트워크 수준 격리
- Key Vault: 비밀 키 중앙 관리"""
            },
            {
                "heading": "네트워크 보안",
                "content": """네트워크 수준의 보안 설정:

NSG (Network Security Group):
- 필요한 포트만 허용 (443, 3443 등)
- 기본 거부 규칙 적용
- Azure 서비스 태그 활용

Private Endpoint:
- 모든 Azure 서비스에 Private Endpoint 적용
- Public 네트워크 접근 비활성화
- Private DNS Zone으로 이름 해석

VNet Peering:
- Hub-Spoke 토폴로지 구성
- Jumpbox를 통한 관리 접근
- Azure Bastion으로 RDP/SSH 보안"""
            },
            {
                "heading": "데이터 보호",
                "content": """데이터 보호를 위한 설정:

암호화:
- Storage: Microsoft 관리 키 또는 CMK
- Key Vault: HSM 보호 키
- 전송 중 암호화: TLS 1.2+

접근 제어:
- Storage Account: 공용 접근 비활성화
- RBAC: Storage Blob Data Reader/Contributor
- SAS 토큰: 최소 유효 기간, IP 제한

로깅 및 모니터링:
- Application Insights: 애플리케이션 로그
- Log Analytics: 리소스 로그
- Azure Monitor: 메트릭 및 알림"""
            }
        ]
    }
}

PRESENTATIONS = {
    "AI_Foundry_아키텍처.pptx": {
        "title": "Azure AI Foundry 아키텍처",
        "category": "아키텍처",
        "slides": [
            {
                "title": "Azure AI Foundry 프라이빗 네트워킹 아키텍처",
                "content": [
                    "East US: AI Foundry Hub, OpenAI, Storage, Key Vault",
                    "Korea Central: Jumpbox VMs, Azure Bastion",
                    "VNet Peering으로 리전 간 연결",
                    "모든 서비스에 Private Endpoint 적용"
                ]
            },
            {
                "title": "네트워크 토폴로지",
                "content": [
                    "vnet-aifoundry (10.0.0.0/16) - East US",
                    "  └ snet-aifoundry (10.0.1.0/24) - Private Endpoints",
                    "vnet-jumpbox-krc (10.1.0.0/16) - Korea Central",
                    "  └ snet-jumpbox (10.1.1.0/24) - Jumpbox VMs",
                    "  └ AzureBastionSubnet (10.1.255.0/26)"
                ]
            },
            {
                "title": "Private Endpoint 구성",
                "content": [
                    "pe-aihub: AI Foundry Hub 접근",
                    "pe-storage-blob, pe-storage-file: 스토리지 접근",
                    "pe-keyvault: Key Vault 접근",
                    "pe-openai: Azure OpenAI 접근",
                    "pe-search: AI Search 접근",
                    "pe-acr: Container Registry 접근"
                ]
            },
            {
                "title": "RAG 패턴 데이터 흐름",
                "content": [
                    "1. 사용자 → Jumpbox → AI Foundry Playground",
                    "2. 질문 입력 → Azure OpenAI",
                    "3. 벡터 검색 → AI Search 인덱스",
                    "4. 관련 문서 검색 → Blob Storage",
                    "5. 컨텍스트 기반 응답 생성"
                ]
            },
            {
                "title": "보안 레이어",
                "content": [
                    "레이어 1: Azure Bastion (외부 접근점)",
                    "레이어 2: NSG (네트워크 필터링)",
                    "레이어 3: Private Endpoint (서비스 격리)",
                    "레이어 4: RBAC (역할 기반 접근 제어)",
                    "레이어 5: Managed Identity (서비스 인증)"
                ]
            }
        ]
    },
    "개발자_온보딩.pptx": {
        "title": "AI Foundry 개발자 온보딩 가이드",
        "category": "온보딩",
        "slides": [
            {
                "title": "환영합니다! AI Foundry 개발 환경",
                "content": [
                    "Azure AI Foundry: 엔터프라이즈 AI 개발 플랫폼",
                    "프라이빗 네트워킹으로 보안 강화",
                    "GPT-4o, AI Search 등 최신 AI 서비스 사용 가능"
                ]
            },
            {
                "title": "접속 방법",
                "content": [
                    "1. Azure Portal 로그인",
                    "2. Bastion 서비스로 Jumpbox 접속",
                    "3. Jumpbox에서 AI Foundry Portal 접근",
                    "Windows Jumpbox IP: 10.1.1.4",
                    "Linux Jumpbox IP: 10.1.1.5"
                ]
            },
            {
                "title": "첫 번째 프로젝트 시작하기",
                "content": [
                    "1. AI Foundry Portal (ai.azure.com) 접속",
                    "2. 프로젝트 선택: aiproj-agents",
                    "3. Playground에서 모델 테스트",
                    "4. 'Add your data'로 RAG 활성화"
                ]
            },
            {
                "title": "사용 가능한 모델",
                "content": [
                    "GPT-4o (gpt-4o): 최신 멀티모달 모델",
                    "  - 2024-11-20 버전 배포됨",
                    "  - 텍스트 + 이미지 처리 가능",
                    "Text Embedding Ada-002:",
                    "  - 벡터 임베딩 생성",
                    "  - RAG 패턴에 필수"
                ]
            },
            {
                "title": "도움이 필요하면?",
                "content": [
                    "📧 AI Platform Team에 문의",
                    "📚 문서: docs 컨테이너 참조",
                    "💡 AI Search 인덱스에서 검색",
                    "🔧 Terraform 코드: infra/ 디렉토리"
                ]
            }
        ]
    }
}


# =============================================================================
# 문서 생성 함수
# =============================================================================
def create_docx(filename: str, doc_info: dict, output_dir: Path) -> Path:
    """DOCX 문서 생성"""
    doc = Document()
    
    # 제목
    title = doc.add_heading(doc_info["title"], 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 메타데이터
    doc.add_paragraph(f"카테고리: {doc_info['category']}")
    doc.add_paragraph(f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    doc.add_paragraph("─" * 50)
    
    # 섹션별 내용
    for section in doc_info["sections"]:
        doc.add_heading(section["heading"], 1)
        for para in section["content"].strip().split("\n\n"):
            doc.add_paragraph(para.strip())
    
    # 저장
    filepath = output_dir / filename
    doc.save(filepath)
    print(f"  ✅ 생성: {filepath}")
    return filepath


def create_pptx(filename: str, ppt_info: dict, output_dir: Path) -> Path:
    """PPTX 프레젠테이션 생성"""
    prs = Presentation()
    
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = ppt_info["title"]
    subtitle.text = f"카테고리: {ppt_info['category']}\n생성일: {datetime.now().strftime('%Y-%m-%d')}"
    
    # 콘텐츠 슬라이드
    bullet_slide_layout = prs.slide_layouts[1]
    
    for slide_info in ppt_info["slides"]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        
        title.text = slide_info["title"]
        tf = body.text_frame
        tf.clear()
        
        for i, item in enumerate(slide_info["content"]):
            if i == 0:
                tf.paragraphs[0].text = item
            else:
                p = tf.add_paragraph()
                p.text = item
                p.level = 1 if item.startswith("  ") else 0
    
    # 저장
    filepath = output_dir / filename
    prs.save(filepath)
    print(f"  ✅ 생성: {filepath}")
    return filepath


# =============================================================================
# Azure 업로드 함수
# =============================================================================
def upload_to_blob(files: list, storage_account: str, container_name: str):
    """Blob Storage에 파일 업로드"""
    print(f"\n📤 Blob Storage 업로드 중...")
    
    credential = DefaultAzureCredential()
    account_url = f"https://{storage_account}.blob.core.windows.net"
    
    try:
        blob_service_client = BlobServiceClient(account_url, credential=credential)
        
        # 컨테이너 생성 (존재하지 않으면)
        container_client = blob_service_client.get_container_client(container_name)
        try:
            container_client.create_container()
            print(f"  ✅ 컨테이너 생성: {container_name}")
        except Exception as e:
            if "ContainerAlreadyExists" in str(e):
                print(f"  ℹ️ 컨테이너 존재: {container_name}")
            else:
                print(f"  ⚠️ 컨테이너 생성 오류: {e}")
        
        # 파일 업로드
        for file_path in files:
            blob_name = file_path.name
            blob_client = container_client.get_blob_client(blob_name)
            
            with open(file_path, "rb") as data:
                blob_client.upload_blob(data, overwrite=True)
            print(f"  ✅ 업로드: {blob_name}")
        
        print(f"\n📁 업로드 완료: {len(files)}개 파일")
        
    except Exception as e:
        print(f"\n❌ Blob 업로드 실패: {e}")
        print("   프라이빗 네트워크 환경에서는 Jumpbox에서 실행해야 합니다.")
        return False
    
    return True


def create_search_index(search_service: str, index_name: str):
    """AI Search 인덱스 생성"""
    print(f"\n🔍 AI Search 인덱스 생성 중...")
    
    credential = DefaultAzureCredential()
    endpoint = f"https://{search_service}.search.windows.net"
    
    try:
        index_client = SearchIndexClient(endpoint, credential)
        
        # 인덱스 스키마 정의
        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True, filterable=True),
            SearchableField(name="content", type=SearchFieldDataType.String, analyzer_name="ko.microsoft"),
            SearchableField(name="title", type=SearchFieldDataType.String, filterable=True, sortable=True),
            SearchableField(name="category", type=SearchFieldDataType.String, filterable=True, facetable=True),
            SimpleField(name="source", type=SearchFieldDataType.String, filterable=True),
            SimpleField(name="metadata_storage_path", type=SearchFieldDataType.String, filterable=True),
            SearchableField(name="metadata_storage_name", type=SearchFieldDataType.String, filterable=True),
        ]
        
        # 벡터 검색 설정
        vector_search = VectorSearch(
            algorithms=[
                HnswAlgorithmConfiguration(name="hnsw-algorithm"),
            ],
            profiles=[
                VectorSearchProfile(name="vector-profile", algorithm_configuration_name="hnsw-algorithm"),
            ],
        )
        
        # Semantic 검색 설정
        semantic_config = SemanticConfiguration(
            name="semantic-config",
            prioritized_fields=SemanticPrioritizedFields(
                content_fields=[SemanticField(field_name="content")],
                title_field=SemanticField(field_name="title"),
            ),
        )
        
        semantic_search = SemanticSearch(configurations=[semantic_config])
        
        # 인덱스 생성
        index = SearchIndex(
            name=index_name,
            fields=fields,
            vector_search=vector_search,
            semantic_search=semantic_search,
        )
        
        result = index_client.create_or_update_index(index)
        print(f"  ✅ 인덱스 생성/업데이트: {result.name}")
        
    except Exception as e:
        print(f"\n❌ 인덱스 생성 실패: {e}")
        print("   프라이빗 네트워크 환경에서는 Jumpbox에서 실행해야 합니다.")
        return False
    
    return True


def index_documents(search_service: str, index_name: str, files: list, doc_infos: dict):
    """문서 인덱싱"""
    print(f"\n📝 문서 인덱싱 중...")
    
    credential = DefaultAzureCredential()
    endpoint = f"https://{search_service}.search.windows.net"
    
    try:
        search_client = SearchClient(endpoint, index_name, credential)
        
        documents = []
        for file_path in files:
            filename = file_path.name
            
            # 문서 정보 찾기
            if filename in doc_infos:
                info = doc_infos[filename]
                # 모든 섹션 콘텐츠 결합
                if "sections" in info:
                    content = "\n\n".join([s["content"] for s in info["sections"]])
                elif "slides" in info:
                    content = "\n\n".join([
                        f"{s['title']}\n" + "\n".join(s["content"])
                        for s in info["slides"]
                    ])
                else:
                    content = ""
                
                doc = {
                    "id": filename.replace(".", "_").replace(" ", "_"),
                    "title": info.get("title", filename),
                    "category": info.get("category", "기타"),
                    "content": content,
                    "source": f"blob://{STORAGE_ACCOUNT}/{CONTAINER_NAME}/{filename}",
                    "metadata_storage_path": f"https://{STORAGE_ACCOUNT}.blob.core.windows.net/{CONTAINER_NAME}/{filename}",
                    "metadata_storage_name": filename,
                }
                documents.append(doc)
        
        # 문서 업로드
        result = search_client.upload_documents(documents)
        success_count = sum(1 for r in result if r.succeeded)
        print(f"  ✅ 인덱싱 완료: {success_count}/{len(documents)}개 문서")
        
    except Exception as e:
        print(f"\n❌ 인덱싱 실패: {e}")
        return False
    
    return True


# =============================================================================
# 메인 실행
# =============================================================================
def main():
    print("=" * 60)
    print("  AI Search RAG 테스트 데이터 생성 및 설정")
    print("=" * 60)
    
    # 출력 디렉토리 생성
    OUTPUT_DIR.mkdir(exist_ok=True)
    print(f"\n📂 출력 디렉토리: {OUTPUT_DIR.absolute()}")
    
    # DOCX 문서 생성
    print("\n📄 DOCX 문서 생성 중...")
    docx_files = []
    for filename, info in DOCUMENTS.items():
        path = create_docx(filename, info, OUTPUT_DIR)
        docx_files.append(path)
    
    # PPTX 프레젠테이션 생성
    print("\n📊 PPTX 프레젠테이션 생성 중...")
    pptx_files = []
    for filename, info in PRESENTATIONS.items():
        path = create_pptx(filename, info, OUTPUT_DIR)
        pptx_files.append(path)
    
    all_files = docx_files + pptx_files
    print(f"\n✅ 총 {len(all_files)}개 문서 생성 완료")
    
    # Azure 리소스에 업로드 (옵션)
    print("\n" + "=" * 60)
    print("  Azure 리소스 설정")
    print("=" * 60)
    
    upload_choice = input("\n📤 Azure Blob Storage에 업로드하시겠습니까? (y/N): ").strip().lower()
    
    if upload_choice == 'y':
        # Blob 업로드
        if upload_to_blob(all_files, STORAGE_ACCOUNT, CONTAINER_NAME):
            # 인덱스 생성
            index_choice = input("\n🔍 AI Search 인덱스를 생성하시겠습니까? (y/N): ").strip().lower()
            if index_choice == 'y':
                if create_search_index(SEARCH_SERVICE, INDEX_NAME):
                    # 문서 인덱싱
                    all_doc_infos = {**DOCUMENTS, **PRESENTATIONS}
                    index_documents(SEARCH_SERVICE, INDEX_NAME, all_files, all_doc_infos)
    
    # 완료 메시지
    print("\n" + "=" * 60)
    print("  설정 완료")
    print("=" * 60)
    print(f"""
📌 생성된 파일:
   {OUTPUT_DIR.absolute()}

📌 AI Foundry에서 RAG 사용 방법:
   1. Jumpbox를 통해 AI Foundry Portal 접속
   2. Project (aiproj-agents) 선택
   3. Playground > 'Add your data' 클릭
   4. Azure AI Search 선택
   5. 인덱스: {INDEX_NAME}
   6. Semantic search 활성화

📌 테스트 질문 예시:
   - "Azure AI Foundry의 프라이빗 네트워킹 구성 방법은?"
   - "RAG 패턴에서 AI Search는 어떻게 사용되나요?"
   - "제로 트러스트 보안 원칙이란?"
""")


if __name__ == "__main__":
    main()
